#Daily News Grabber
#Harris Ransom

#Daily News Grabber
#Harris Ransom
#Libraries Used:
#   -python-pptx
#   -feedparser

#Imports libraries
import feedparser
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

#Creates lists to hold headlines and links
News = []
Links = []

#Creates news feed and pulls down latest news
SKYnews = feedparser.parse("http://feeds.skynews.com/feeds/rss/uk.xml")
#BBCnews = feedparser.parse("http://feeds.bbci.co.uk/news/rss.xml?edition=uk")
#NYTnews = feedparser.parse("https://archive.nytimes.com/www.nytimes.com/services/xml/rss/index.html?mcubz=0")
#WPnews = feedparser.parse("https://www.washingtonpost.com/discussions/2018/10/12/washington-post-rss-feeds/?noredirect=on")

#Creates the slide layout
def create_my_default_slide(title, subtitle):
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    slide.shapes.title.text = title
    slide.placeholders[1].text = subtitle
    return slide

#Creates the hyperlink layout
def add_hyper_link(shape, text, url):
    p = shape.text_frame.paragraphs[0]
    r = p.add_run()
    r.text = text
    p.font.size = Pt(12)
    p.font.name = 'Lato'
    p.font.color.rgb = RGBColor(0, 0, 255)
    hlink = r.hyperlink
    hlink.address = url

#Pulls top 10 news stories
for i in range(10):
    text = SKYnews["entries"][i]["title"]
    links = SKYnews["entries"][i]["link"]
    media = SKYnews["entries"][i]["media_thumbnail"]
    print(text)
    print(links)
    print(media)
    News.append(text) #add headlines to News list
    Links.append(links) # add links to Links list

#Creates a dictionary from the News and Links lists
News_Dict = {}
for i in range(len(News)):
    News_Dict [News[i]] = Links[i]
print ("")   
print(News_Dict) 

#Creates the powerpoint
prs = Presentation()
for key, value in News_Dict.items():
    this_slide = create_my_default_slide("%s" % key, "Click for full story: ")
    add_hyper_link(this_slide.shapes[1], value, value)

print(len(prs.slides))

#Saves the powerpoint
prs.save('dailyNews.pptx')